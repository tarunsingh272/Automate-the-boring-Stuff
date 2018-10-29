# -*- coding: utf-8 -*-
"""
Created on Wed Aug  1 14:05:24 2018

@author: henri
"""
from pptx import Presentation

def update_pres_text(template_name):
#    from pptx import Presentation
    import pandas as pd
    import numpy as np
    
    week_num = int(input('enter week to analyze: '))
    
    xls = pd.ExcelFile('graphs_to_slides2.xlsx')
    df1 = xls.parse(0)
    df2 = xls.parse(1, index_col=0)
    df3 = xls.parse(2, index_col=0)
    
    # apply multi column index to df4
    df4 = xls.parse(3, header = 1)
    groups = int(len(df4.columns) / 4)
    df4.columns = [['week' + str(week_num + i) for i in range(groups) for j in range(4)],list(df4.columns)]
    
    # create the text variables
    X1 = week_num  # assign the report week
    X2 = df1[df1['Week'] == week_num]['Hours Booked'].values[0] # hours booked for specific week
    X3 = df1[df1['Week'] == week_num]['Hours Expected'].values[0]  # hours expected for specific week
    X4 = df1[df1['Week'] == week_num]['Lag/Lead'].values[0]  # YTD target lag/lead hours
    X5 = (df1[df1['Week'] == week_num]['Hours/Resource'].values[0] / 40) * 100  # resource utilization
    X6 = "NA"
    X7 = "NA"
    X8 = "NA"
    X9 = "NA"
    df2['util'] = df2.sum(axis = 1)  # create max utilization column
    X10 = df2[df2['util'] == max(df2['util'])].index.values[0]  # name of project with max 3 week utilization
    X11 = df2[df2['util'] == min(df2['util'])].index.values[0]  # name of project with min 3 week utilization
    X12 = df2.loc[X10][2]/df2.loc[X10][3] * 100  # print the change in utiliztion for the max project
    X13 = df2.loc[X11][2]/df2.loc[X11][3] * 100  # print the change in utiliztion for the max project
    X14 = df3[df3['Growth'] == max(df3['Growth'])].index.values[0]  # find department with max growth
    X15 = df3[df3['Growth'] == min(df3['Growth'])].index.values[0]  # find department with min growth
    
    # add a row that is the ratio of dev to maintenance
    df4 = df4.append(pd.Series(df4.loc['Development '] / df4.loc['Maintenance'], name='ratio'))
    X16 = round(df4.loc['ratio'].mean(), 2)  # average ratio across three weeks as integer

    values = {'X1':X1, 'X2':X2, 'X3':X3, 'X4':X4, 'X5':X5, 
              'X6':X6, 'X7':X7, 'X8':X8, 'X9':X9, 'X10':X10, 
              'X11':X11, 'X12':X12, 'X13':X13, 'X14':X14, 
              'X15':X15, 'X16':X16}
    
    prs = Presentation(template_name)
    
    for i in range(16, 0, -1):
        key = 'X' + str(i)
        replaceText(key, values[key], prs)
    
    prs.save('deck_format2.pptx')

def replaceText(key, value, prs):
    
    if type(value) == int:
        formatted_value = str(value)
    elif type(value) != str:
        formatted_value = format(float(value), ',.1f')
    else:
        formatted_value = str(value)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                location = shape.text.find(key)
                while location != -1:
                    old_text = shape.text
                    new_text = formatted_value
                    shape.text = old_text[:location] + new_text + old_text[location + len(key):]
                    location = shape.text.find(key)

update_pres_text('deck_format.pptx')
    
